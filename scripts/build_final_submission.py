"""Update Word Report & PowerPoint with student details"""
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PPRGB
from pptx.enum.text import PP_ALIGN

# ── Student Details ───────────────────────────────────────────────────────────
NAME     = "Rahul Parihar"
ROLL     = "268"
COLLEGE  = "JECRC Foundation, Jaipur"
PROF     = "Prof. Ram Singh"
SUBJECT  = "Business & Economic Environment (BEE) Lab"
DATE     = "April 2026"
DEPT     = "B.Tech First Year"

DB = RGBColor(0x1F,0x4E,0x79)
WT = RGBColor(0xFF,0xFF,0xFF)
DG = RGBColor(0x33,0x33,0x33)
MB = RGBColor(0x2E,0x75,0xB6)

def shd(cell, hex_c):
    tc=cell._tc; p=tc.get_or_add_tcPr(); s=OxmlElement('w:shd')
    s.set(qn('w:val'),'clear'); s.set(qn('w:color'),'auto')
    s.set(qn('w:fill'),hex_c); p.append(s)

def fnt(bold=False, sz=11, color=DG):
    from docx.shared import Pt
    f_obj = type('F', (), {})()
    return (bold, sz, color)

def set_cell(cell, text, bg_hex, txt_color, bold=False, sz=10.5, align=WD_ALIGN_PARAGRAPH.LEFT):
    shd(cell, bg_hex)
    cell.paragraphs[0].clear()
    p = cell.paragraphs[0]; p.alignment = align
    r = p.add_run(text); r.font.bold = bold
    r.font.size = Pt(sz); r.font.name = "Calibri"
    r.font.color.rgb = txt_color

# ═══════════════════════════════════════════════════════════════
# BUILD WORD REPORT FROM SCRATCH WITH CORRECT DETAILS
# ═══════════════════════════════════════════════════════════════
from docx.shared import Inches as WInches

doc = Document()
for sec in doc.sections:
    sec.top_margin=Cm(2); sec.bottom_margin=Cm(2)
    sec.left_margin=Cm(2.5); sec.right_margin=Cm(2)

def hd(text, lvl=1):
    p = doc.add_heading(text, level=lvl)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    clr = DB if lvl==1 else MB
    for r in p.runs:
        r.font.color.rgb = clr; r.font.name = "Calibri"
        r.font.size = Pt(16 if lvl==1 else 13)

def para(text, sz=11, bold=False, italic=False, align=WD_ALIGN_PARAGRAPH.LEFT, space=6):
    p = doc.add_paragraph(); p.alignment = align
    p.paragraph_format.space_after = Pt(space)
    r = p.add_run(text); r.bold=bold; r.italic=italic
    r.font.size=Pt(sz); r.font.name="Calibri"; r.font.color.rgb=DG

def bul(text):
    p = doc.add_paragraph(style='List Bullet')
    p.paragraph_format.space_after = Pt(3)
    r = p.add_run(text); r.font.size=Pt(11); r.font.name="Calibri"; r.font.color.rgb=DG

def make_table(cols, data):
    t = doc.add_table(rows=len(data)+1, cols=len(cols))
    t.style = 'Table Grid'
    # header
    for i, c in enumerate(cols):
        cell = t.rows[0].cells[i]
        set_cell(cell, c, "1F4E79", WT, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    # data rows
    for ri, row in enumerate(data, 1):
        bg = "EBF3FB" if ri%2==0 else "FFFFFF"
        for ci, val in enumerate(row):
            set_cell(t.rows[ri].cells[ci], str(val), bg, DG, sz=10)
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

# ── COVER PAGE ───────────────────────────────────────────────────────────────
para("", sz=2)
p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run("BEE LAB PROJECT REPORT")
r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=DB; r.font.name="Calibri"

p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run("Retail Sales Analytics: Excel & Power BI Dashboards")
r.font.size=Pt(15); r.font.color.rgb=MB; r.font.name="Calibri"

para("", sz=5)

# Cover table
ct = doc.add_table(rows=6, cols=2); ct.style='Table Grid'
cover = [
    ("Student Name",  NAME),
    ("Roll Number",   ROLL),
    ("College",       COLLEGE),
    ("Department",    DEPT),
    ("Faculty Guide", PROF),
    ("Date",          DATE),
]
for i,(lbl,val) in enumerate(cover):
    set_cell(ct.rows[i].cells[0], lbl, "1F4E79", WT, bold=True)
    set_cell(ct.rows[i].cells[1], val, "EBF3FB" if i%2==0 else "FFFFFF", DG)

doc.add_page_break()

# ── 1. INTRODUCTION ──────────────────────────────────────────────────────────
hd("1. Introduction & Background")
para("Retail sales data contains rich insights about customer behaviour, product performance, and regional trends. This project uses Microsoft Excel and Microsoft Power BI to transform a 300-order, 3-year (2020–2022) retail dataset into two interactive Business Intelligence dashboards that help stakeholders make fast, data-driven decisions without requiring any programming knowledge.")

# ── 2. PROBLEM STATEMENT ─────────────────────────────────────────────────────
hd("2. Problem Statement")
p = doc.add_paragraph(); p.paragraph_format.space_after=Pt(6)
r = p.add_run('"Retail businesses generate large volumes of transactional data but often lack the tools to extract actionable insights. This project designs two interactive dashboards — one in Microsoft Excel and one in Microsoft Power BI — enabling stakeholders to monitor KPIs, compare category and regional performance, analyze profit margins, and forecast future sales trends without programming knowledge."')
r.italic=True; r.font.size=Pt(11); r.font.name="Calibri"; r.font.color.rgb=DG

# ── 3. OBJECTIVES ────────────────────────────────────────────────────────────
hd("3. Objectives")
for o in ["Build an interactive Excel Dashboard with Pivot Tables, PivotCharts, and Slicers.",
          "Develop a Power BI Dashboard with 16 DAX-powered KPI measures and a 3-month forecast.",
          "Perform Year-over-Year (YoY) and Month-over-Month (MoM) trend analysis.",
          "Identify top-performing and loss-making product sub-categories.",
          "Deliver a professional report and presentation for submission."]:
    bul(o)

# ── 4. METHODOLOGY ───────────────────────────────────────────────────────────
hd("4. Methodology")
make_table(["Phase","Tool","Activity"],[
    ("1 — Data Collection","Kaggle / CSV","Downloaded retail sales dataset (300 rows, 2020-2022)"),
    ("2 — Data Cleaning","Power Query","Fixed data types, standardized dates, removed nulls"),
    ("3 — Data Modeling","Power BI DAX","Created DateTable, Star Schema, 16 DAX measures"),
    ("4 — Dashboard Design","Excel + Power BI","Pivot Tables, PivotCharts, Slicers, Visuals, Forecast"),
    ("5 — Analysis & Report","Both Tools","KPI tracking, YoY trends, findings, submission"),
])

# ── 5. DATASET ───────────────────────────────────────────────────────────────
hd("5. Dataset Description")
para("File: retail_sales_final.csv  |  Source: Kaggle Superstore (curated)", sz=11)
for s in ["Total Records: 300 orders across 3 years (2020, 2021, 2022)",
          "Categories: Technology ($2,50,980) | Furniture ($74,164) | Office Supplies ($9,251)",
          "Regions: West ($1,01,195) | South ($84,679) | East ($81,436) | Central ($67,086)",
          "Segments: Consumer (48%) | Corporate (30%) | Home Office (22%)",
          "21 loss-making orders due to high discounts on Furniture"]:
    bul(s)

# ── 6. DATA DICTIONARY ───────────────────────────────────────────────────────
hd("6. Data Dictionary")
make_table(["Column","Type","Description"],[
    ("Order ID","Text","Unique order identifier (CA-YYYY-XXXXXX)"),
    ("Order Date","Date","Date order was placed (DD-MM-YYYY)"),
    ("Ship Date","Date","Date order was shipped (DD-MM-YYYY)"),
    ("Customer Name","Text","Full name of customer"),
    ("Segment","Text","Consumer / Corporate / Home Office"),
    ("Region","Text","East / West / Central / South"),
    ("State","Text","US delivery state"),
    ("Category","Text","Technology / Furniture / Office Supplies"),
    ("Sub-Category","Text","Phones, Laptops, Chairs, Tables, Binders, etc."),
    ("Product Name","Text","Full product name"),
    ("Sales","Decimal","Gross revenue ($)"),
    ("Quantity","Integer","Units ordered"),
    ("Discount","Decimal","Discount rate (0.10 = 10%)"),
    ("Profit","Decimal","Net profit; negative = loss"),
])

# ── 7. EXCEL DASHBOARD ───────────────────────────────────────────────────────
hd("7. Excel Dashboard — Store Sales Performance")
para("File: Store_Sales_Dashboard_FINAL.xlsx  |  6 sheets with charts and KPI cards", sz=11)
hd("7.1 KPI Summary", lvl=2)
make_table(["KPI Metric","Value"],[
    ("Total Sales (2020–2022)","$3,34,395.01"),
    ("Total Profit","$51,612.01"),
    ("Profit Margin %","15.4%"),
    ("Total Orders","300"),
    ("Total Units Sold","699"),
    ("2020 Sales","$81,018.48"),
    ("2021 Sales","$1,31,145.77  (+61.9% YoY)"),
    ("2022 Sales","$1,22,230.76  (-6.8% YoY)"),
])
hd("7.2 Dashboard Sheets", lvl=2)
make_table(["Sheet Name","Contents"],[
    ("DASHBOARD","10 KPI cards, 3 PivotCharts, segment & year summary"),
    ("SalesData","300-row formatted Excel Table"),
    ("PT_Category","Sales & Profit by Category × Region + Bar Chart"),
    ("PT_YearlyComparison","2020 vs 2021 vs 2022 monthly trend + Line Chart"),
    ("PT_TopProducts","Top 10 products by profit + Horizontal Bar"),
    ("PT_SubCategory","All sub-categories ranked by profit (losses in red)"),
])
hd("7.3 Slicers & Interactivity", lvl=2)
for s in ["Region Slicer — filters all charts simultaneously via Report Connections",
          "Category Slicer — isolates a single product category",
          "Segment Slicer — Consumer / Corporate / Home Office",
          "Year Slicer — compare any single year or multi-year view",
          "Date Timeline Slicer — drag to select custom date range"]:
    bul(s)

# ── 8. POWER BI DASHBOARD ────────────────────────────────────────────────────
hd("8. Power BI Dashboard — Retail Sales & Forecast BI")
hd("8.1 Data Model (Star Schema)", lvl=2)
for s in ["SalesData (Fact Table) — 300 rows, 14 columns",
          "DateTable (Dimension) — created via DAX CALENDAR() function",
          "Relationship: DateTable[Date]  ──1:Many──  SalesData[Order Date]",
          "DateTable marked as official Date Table for all time intelligence measures"]:
    bul(s)
hd("8.2 Key DAX Measures", lvl=2)
make_table(["Measure","Formula (simplified)","Result"],[
    ("Total Sales","SUM(SalesData[Sales])","$3,34,395"),
    ("Total Profit","SUM(SalesData[Profit])","$51,612"),
    ("Profit Margin %","DIVIDE([Profit],[Sales],0)","15.4%"),
    ("YTD Sales","TOTALYTD([Total Sales],DateTable[Date])","Cumulative"),
    ("YoY Growth %","DIVIDE([Sales]-[Sales PY],[Sales PY],0)","+61.9% (2021)"),
    ("MoM Growth %","DIVIDE([Sales]-[Sales PM],[Sales PM],0)","Monthly"),
    ("Sales Previous Year","CALCULATE([Sales],SAMEPERIODLASTYEAR(...))","Benchmark"),
    ("Loss Orders Count","CALCULATE([Orders],Profit<0)","21 orders"),
])
hd("8.3 Report Pages", lvl=2)
make_table(["Page","Visuals"],[
    ("Overview","5 KPI Cards, Bar Chart, Line Chart, Donut Chart, Map Visual"),
    ("Category Analysis","Matrix Table, Stacked Bar, Scatter Plot, KPI Card"),
    ("Regional View","Filled Map, Column Chart, Summary Table, Treemap"),
    ("Forecast","Line Chart + 3-month forecast (95% CI), MTD Sales, MoM Cards"),
])
hd("8.4 Forecasting Setup", lvl=2)
for s in ["Visual: Line Chart — X-Axis: DateTable[Date] (Month), Y-Axis: Total Sales",
          "Analytics Pane → Forecast → Add",
          "Forecast length: 3 months  |  Confidence Interval: 95%  |  Seasonality: Auto",
          "Output: Dashed forecast line with grey confidence band"]:
    bul(s)

# ── 9. KEY FINDINGS ──────────────────────────────────────────────────────────
hd("9. Key Findings & Business Insights")
for f in ["Technology is the top category at $2,50,980 revenue (75% of total) — Laptops & Phones lead.",
          "2021 recorded the highest YoY growth at +61.9% — strong post-2020 recovery.",
          "2022 saw a -6.8% decline — warrants investigation into pricing/discount strategy.",
          "West region leads all regions with $1,01,195 in total sales.",
          "Laptops ($19,678) and Phones ($19,621) are the most profitable sub-categories.",
          "Tables and Bookcases show losses due to heavy discounts (15–20%).",
          "Consumer segment contributes 48% of total revenue — the primary market.",
          "Forecast: Sales projected to stabilize and recover in early 2023."]:
    bul(f)

# ── 10. LIMITATIONS ──────────────────────────────────────────────────────────
hd("10. Limitations")
for l in ["Dataset is US-based; Indian retail patterns may differ significantly.",
          "Power BI forecasting uses exponential smoothing only — no external factor modeling.",
          "300 records is sufficient for learning purposes; enterprise dashboards use millions of rows.",
          "No live database connection — CSV is static and not updated in real-time.",
          "No customer-level analytics (RFM, CLV) — out of scope for this project."]:
    bul(l)

# ── 11. CONCLUSION ───────────────────────────────────────────────────────────
hd("11. Conclusion")
para("This project successfully demonstrates a complete Business Intelligence workflow — from raw retail CSV data to two fully interactive dashboards using Microsoft Excel and Power BI. The Excel dashboard provides Pivot Table-based KPI monitoring with Slicer interactivity, while the Power BI dashboard delivers advanced DAX analytics, a proper Star Schema data model, and 3-month sales forecasting with 95% confidence intervals. Both tools together address all project objectives of the BEE Lab and provide a strong foundation for enterprise-level retail analytics.")

# ── 12. REFERENCES ───────────────────────────────────────────────────────────
hd("12. References")
for ref in ["[1] Microsoft DAX Reference — docs.microsoft.com/en-us/dax/",
            "[2] Kaggle Superstore Dataset — kaggle.com/datasets/vivek468/superstore-dataset-final",
            "[3] Power BI Forecasting — learn.microsoft.com/en-us/power-bi/visuals",
            "[4] Excel PivotTable Guide — support.microsoft.com/en-us/office",
            "[5] Power BI Desktop Download — powerbi.microsoft.com/desktop",
            "[6] openpyxl Documentation — openpyxl.readthedocs.io"]:
    bul(ref)

WORD_OUT = r"e:\BEE\Section4_Report\BEE_Lab_Report_RAHUL_PARIHAR.docx"
doc.save(WORD_OUT)
print(f"Word Report saved -> {WORD_OUT}")

# ═══════════════════════════════════════════════════════════════
# UPDATE POWERPOINT WITH STUDENT DETAILS
# ═══════════════════════════════════════════════════════════════
prs = Presentation(r"e:\BEE\Section4_Report\BEE_Lab_Presentation_FINAL.pptx")

def upd_text(shape, old, new):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

replacements = {
    "Group 4  |  B.Tech First Year  |  Business & Economic Environment (BEE) Lab  |  April 2026":
        f"{NAME}  |  Roll No: {ROLL}  |  {COLLEGE}  |  {SUBJECT}  |  {DATE}",
    "Group 4  |  B.Tech First Year  |  BEE Lab  |  April 2026":
        f"{NAME}  |  Roll No: {ROLL}  |  {COLLEGE}  |  {DATE}",
}

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                for old, new in replacements.items():
                    if old in full:
                        if para.runs:
                            para.runs[0].text = new
                            for r in para.runs[1:]: r.text = ""

# Fix slide 1 — add proper student info box
sl1 = prs.slides[0]
from pptx.util import Inches, Pt as PPt
txb = sl1.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12), Inches(0.5))
tf = txb.text_frame; p2 = tf.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = f"Submitted by: {NAME}  |  Roll No: {ROLL}  |  {COLLEGE}"
r2.font.size = PPt(13); r2.font.bold = True
r2.font.color.rgb = PPRGB(0xBD,0xD7,0xEE); r2.font.name = "Calibri"

txb2 = sl1.shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(12), Inches(0.4))
tf2 = txb2.text_frame; p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = f"Faculty Guide: {PROF}  |  {SUBJECT}"
r3.font.size = PPt(12); r3.font.color.rgb = PPRGB(0xBD,0xD7,0xEE); r3.font.name = "Calibri"

PPTX_OUT = r"e:\BEE\Section4_Report\BEE_Lab_Presentation_RAHUL_PARIHAR.pptx"
prs.save(PPTX_OUT)
print(f"PowerPoint saved -> {PPTX_OUT}")

# ── Final ZIP ─────────────────────────────────────────────────────────────────
import zipfile, os
files_to_zip = [
    r"e:\BEE\Section1_Dataset\retail_sales_final.csv",
    r"e:\BEE\Section2_Excel\Store_Sales_Dashboard_FINAL.xlsx",
    r"e:\BEE\Section3_PowerBI\DAX_Measures.dax",
    r"e:\BEE\Section3_PowerBI\PowerBI_Setup_Guide.md",
    WORD_OUT,
    PPTX_OUT,
    r"e:\BEE\README.md",
]
ZIP_OUT = r"e:\BEE\BEE_Lab_RAHUL_PARIHAR_268.zip"
with zipfile.ZipFile(ZIP_OUT, 'w', zipfile.ZIP_DEFLATED) as zf:
    for fp in files_to_zip:
        zf.write(fp, os.path.basename(fp))
print(f"ZIP saved -> {ZIP_OUT}")
print(f"ZIP size : {os.path.getsize(ZIP_OUT)/1024:.1f} KB")
print()
print("=" * 55)
print(f"  Student  : {NAME}  |  Roll: {ROLL}")
print(f"  College  : {COLLEGE}")
print(f"  Faculty  : {PROF}")
print(f"  Files    : Word + PPTX + Excel + CSV + DAX + Guide")
print("=" * 55)
