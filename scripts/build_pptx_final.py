"""PowerPoint Presentation — BEE Lab Group 4"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

DB = RGBColor(0x1F,0x4E,0x79); MB = RGBColor(0x2E,0x75,0xB6)
WT = RGBColor(0xFF,0xFF,0xFF); DG = RGBColor(0x33,0x33,0x33)
GR = RGBColor(0x37,0x56,0x23); OR = RGBColor(0xC5,0x5A,0x11)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, hex_color):
    fill = slide.background.fill; fill.solid()
    fill.fore_color.rgb = RGBColor(int(hex_color[0:2],16),int(hex_color[2:4],16),int(hex_color[4:6],16))

def box(slide, l, t, w, h, text, fsize=24, bold=True, color=WT, bg_color=None,
        align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fsize); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = "Calibri"
    if bg_color:
        fill = txBox.fill; fill.solid()
        fill.fore_color.rgb = RGBColor(int(bg_color[0:2],16),int(bg_color[2:4],16),int(bg_color[4:6],16))
    return txBox

def rect(slide, l, t, w, h, hex_color, text="", fsize=18, bold=True, txt_color=WT, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(int(hex_color[0:2],16),int(hex_color[2:4],16),int(hex_color[4:6],16))
    shape.line.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    if text:
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text; r.font.size = Pt(fsize)
        r.font.bold = bold; r.font.color.rgb = txt_color; r.font.name = "Calibri"
    return shape

def line_sep(slide, t, color="2E75B6"):
    rect(slide, 0.3, t, 12.7, 0.04, color)

# ═══ SLIDE 1 — TITLE ══════════════════════════════════════════════════════════
sl1 = blank(); bg(sl1,"1F4E79")
rect(sl1,0,0,13.33,7.5,"1F4E79")
rect(sl1,0,0,13.33,0.15,"2E75B6")
rect(sl1,0,7.35,13.33,0.15,"2E75B6")
box(sl1,0.6,0.8,12,1.5,"BEE Lab Project Report",46,True,WT,align=PP_ALIGN.CENTER)
box(sl1,0.6,2.3,12,0.8,"Retail Sales Analytics: Excel & Power BI Dashboards",24,False,
    RGBColor(0xBD,0xD7,0xEE),align=PP_ALIGN.CENTER)
line_sep(sl1,3.2)
kpis = [("$334,395","Total Sales"),("$51,612","Total Profit"),("15.4%","Profit Margin"),("300","Orders")]
for i,(val,lbl) in enumerate(kpis):
    x = 0.7 + i*3.2
    rect(sl1,x,3.5,2.8,1.1,"2E75B6")
    box(sl1,x+0.1,3.55,2.6,0.5,val,22,True,WT,align=PP_ALIGN.CENTER)
    box(sl1,x+0.1,4.0,2.6,0.4,lbl,12,False,RGBColor(0xBD,0xD7,0xEE),align=PP_ALIGN.CENTER)
box(sl1,0.6,4.9,12,0.5,"Group 4  |  B.Tech First Year  |  Business & Economic Environment (BEE) Lab  |  April 2026",
    14,False,RGBColor(0xBD,0xD7,0xEE),align=PP_ALIGN.CENTER)

# ═══ SLIDE 2 — AGENDA ═════════════════════════════════════════════════════════
sl2 = blank(); bg(sl2,"F4F6FB")
rect(sl2,0,0,13.33,1.1,"1F4E79")
box(sl2,0.5,0.15,12,0.8,"Agenda / Table of Contents",28,True,WT,align=PP_ALIGN.LEFT)
items = ["01  Problem Statement & Objectives",
         "02  Dataset Overview (300 Orders | 2020-2022)",
         "03  Excel Dashboard — Store Sales Performance",
         "04  Power BI Dashboard — Retail Sales & Forecast BI",
         "05  Key DAX Measures Explained",
         "06  Key Findings & Business Insights",
         "07  Limitations & Conclusion"]
colors=["1F4E79","2E75B6","375623","7030A0","C55A11","833C00","1F4E79"]
for i,(item,c) in enumerate(zip(items,colors)):
    rect(sl2,0.5,1.3+i*0.84,12.3,0.7,c,item,16,False,WT,PP_ALIGN.LEFT)

# ═══ SLIDE 3 — PROBLEM & OBJECTIVES ══════════════════════════════════════════
sl3 = blank(); bg(sl3,"F4F6FB")
rect(sl3,0,0,13.33,1.1,"1F4E79")
box(sl3,0.5,0.15,12,0.8,"Problem Statement & Objectives",28,True,WT)
box(sl3,0.5,1.2,12.3,1.5,
    '"Retail businesses generate large volumes of transactional data but often lack tools to extract insights.\n'
    'This project builds two interactive BI dashboards — in Excel and Power BI — enabling stakeholders\n'
    'to monitor KPIs, track trends, and forecast sales without programming knowledge."',
    14,False,DG)
rect(sl3,0.5,2.8,12.3,0.5,"2E75B6","Objectives",16,True,WT)
objs=["Build Excel Dashboard with Pivot Tables, PivotCharts & Slicers",
      "Build Power BI Dashboard with 16 DAX measures & 3-month Forecast",
      "Perform 3-Year YoY & MoM Sales Trend Analysis (2020-2022)",
      "Identify top products, loss-making sub-categories, best regions"]
for i,obj in enumerate(objs):
    rect(sl3,0.5,3.4+i*0.78,12.3,0.68,"EBF3FB","  "+obj,14,False,DG,PP_ALIGN.LEFT)

# ═══ SLIDE 4 — DATASET ════════════════════════════════════════════════════════
sl4 = blank(); bg(sl4,"F4F6FB")
rect(sl4,0,0,13.33,1.1,"1F4E79")
box(sl4,0.5,0.15,12,0.8,"Dataset Overview",28,True,WT)
stats=[("300 Orders","Total Records"),("3 Years","2020 / 2021 / 2022"),
       ("14 Columns","Full Attribute Set"),("3 Categories","Tech | Furniture | Office"),
       ("4 Regions","E | W | Central | S"),("21 Losses","High-discount orders")]
for i,(val,lbl) in enumerate(stats):
    x=0.5+(i%3)*4.2; y=1.3+(i//3)*1.8
    rect(sl4,x,y,3.8,0.7,["1F4E79","2E75B6","375623","7030A0","C55A11","833C00"][i])
    box(sl4,x+0.1,y+0.05,3.6,0.35,val,20,True,WT,align=PP_ALIGN.CENTER)
    box(sl4,x+0.1,y+0.38,3.6,0.25,lbl,12,False,WT,align=PP_ALIGN.CENTER)
rect(sl4,0.5,4.9,12.3,0.5,"2E75B6","File: retail_sales_final.csv  |  Source: Kaggle Superstore Dataset",14,False,WT)
box(sl4,0.5,5.5,12.3,1.6,
    "Category Split:  Technology $250,980 (75%)  |  Furniture $74,164 (22%)  |  Office Supplies $9,251 (3%)\n"
    "Region Split:    West $101,195  |  South $84,679  |  East $81,436  |  Central $67,086\n"
    "Segment Split:   Consumer $159,089 (48%)  |  Corporate $101,812 (30%)  |  Home Office $73,494 (22%)",
    13,False,DG)

# ═══ SLIDE 5 — EXCEL DASHBOARD ════════════════════════════════════════════════
sl5 = blank(); bg(sl5,"F4F6FB")
rect(sl5,0,0,13.33,1.1,"375623")
box(sl5,0.5,0.15,12,0.8,"Section 2: Excel Dashboard — Store Sales Performance",26,True,WT)
sheets=[("DASHBOARD","KPI cards (10 total) + 3 PivotCharts","1F4E79"),
        ("SalesData","300-row formatted Excel Table","2E75B6"),
        ("PT_Category","Sales & Profit by Category × Region + Bar Chart","375623"),
        ("PT_YearlyComparison","2020 vs 2021 vs 2022 monthly trend + Line Chart","7030A0"),
        ("PT_TopProducts","Top 10 products by profit + Horizontal Bar","C55A11"),
        ("PT_SubCategory","All sub-categories ranked by profit (losses in red)","833C00")]
for i,(sh,desc,c) in enumerate(sheets):
    x=0.5+(i%2)*6.4; y=1.3+(i//2)*1.8
    rect(sl5,x,y,5.9,0.5,c,sh,15,True,WT)
    box(sl5,x+0.1,y+0.55,5.7,1.1,desc,13,False,DG)
rect(sl5,0.5,6.9,12.3,0.4,"375623","Slicers: Region | Category | Segment | Year | Date Timeline — all connected via Report Connections",12,False,WT)

# ═══ SLIDE 6 — POWER BI ════════════════════════════════════════════════════════
sl6 = blank(); bg(sl6,"F4F6FB")
rect(sl6,0,0,13.33,1.1,"7030A0")
box(sl6,0.5,0.15,12,0.8,"Section 3: Power BI Dashboard — Retail Sales & Forecast BI",26,True,WT)
rect(sl6,0.5,1.2,5.8,1.0,"1F4E79","Data Model: Star Schema",16,True,WT)
box(sl6,0.5,2.3,5.8,1.4,"DateTable [1] ──────── [*] SalesData\nRelationship: DateTable[Date] → SalesData[Order Date]\nDateTable marked as official Date Table",13,False,DG)
rect(sl6,6.9,1.2,5.9,1.0,"2E75B6","Report Pages (4 total)",16,True,WT)
for i,pg in enumerate(["Overview  — KPI Cards + Map + Charts","Category Analysis  — Matrix + Scatter",
                        "Regional View  — Filled Map + Treemap","Forecast  — 3-month projection (95% CI)"]):
    box(sl6,6.9,2.3+i*0.35,5.9,0.33,f"  {i+1}. {pg}",12,False,DG)
rect(sl6,0.5,4.0,12.3,0.5,"7030A0","Slicers: Date Range | Category | Region | Segment  (synced across all 4 pages)",14,False,WT)
box(sl6,0.5,4.6,12.3,0.5,"Forecasting: Analytics Pane → Forecast → 3 months ahead | Confidence: 95% | Seasonality: Auto",13,False,DG)
rect(sl6,0.5,5.3,12.3,0.5,"375623","Export: File → Export → PDF  |  Save as .pbix for submission",14,False,WT)

# ═══ SLIDE 7 — DAX MEASURES ════════════════════════════════════════════════════
sl7 = blank(); bg(sl7,"F4F6FB")
rect(sl7,0,0,13.33,1.1,"C55A11")
box(sl7,0.5,0.15,12,0.8,"Key DAX Measures — Power BI",28,True,WT)
measures=[("Total Sales","SUM(SalesData[Sales])","$334,395"),
          ("Profit Margin %","DIVIDE([Profit],[Sales],0)","15.4%"),
          ("YTD Sales","TOTALYTD([Total Sales],DateTable[Date])","Cumulative"),
          ("YoY Growth %","DIVIDE([Sales]-[Sales PY],[Sales PY],0)","+61.9%"),
          ("MoM Growth %","DIVIDE([Sales]-[Sales PM],[Sales PM],0)","Monthly"),
          ("Loss Orders","CALCULATE([Orders],Profit<0)","21 orders")]
for i,(name,formula,result) in enumerate(measures):
    x=0.5+(i%2)*6.4; y=1.3+(i//2)*1.8
    rect(sl7,x,y,2.2,1.5,"1F4E79",name,13,True,WT)
    box(sl7,x+2.3,y+0.1,3.0,0.7,formula,11,False,DG)
    rect(sl7,x+2.3,y+0.9,3.0,0.5,"375623",result,13,True,WT)

# ═══ SLIDE 8 — KEY FINDINGS ════════════════════════════════════════════════════
sl8 = blank(); bg(sl8,"F4F6FB")
rect(sl8,0,0,13.33,1.1,"1F4E79")
box(sl8,0.5,0.15,12,0.8,"Key Findings & Business Insights",28,True,WT)
findings=[("Technology #1","75% of revenue — Laptops & Phones lead","1F4E79"),
          ("2021 Best Year","YoY Growth +61.9% — highest across 3 years","375623"),
          ("West Leads","$101,195 in sales — highest region","2E75B6"),
          ("Loss Alert","Bookcases & Tables lose money (high discounts)","C55A11"),
          ("Consumer King","48% revenue share — largest segment","7030A0"),
          ("2022 Dip","-6.8% YoY — investigate cause for recovery plan","833C00")]
for i,(title,desc,c) in enumerate(findings):
    x=0.5+(i%3)*4.2; y=1.3+(i//3)*2.5
    rect(sl8,x,y,3.9,0.6,c,title,16,True,WT)
    box(sl8,x+0.1,y+0.65,3.7,1.5,desc,13,False,DG)

# ═══ SLIDE 9 — CONCLUSION ═════════════════════════════════════════════════════
sl9 = blank(); bg(sl9,"1F4E79")
rect(sl9,0,0,13.33,7.5,"1F4E79")
rect(sl9,0,0,13.33,0.12,"2E75B6"); rect(sl9,0,7.38,13.33,0.12,"2E75B6")
box(sl9,0.6,0.5,12,0.8,"Conclusion",38,True,WT,align=PP_ALIGN.CENTER)
line_sep(sl9,1.4)
box(sl9,0.6,1.6,12,1.8,
    "This project successfully demonstrates a complete Business Intelligence workflow —\n"
    "from raw retail CSV data to two fully interactive dashboards using Excel and Power BI.\n"
    "The Excel dashboard enables Pivot Table-based KPI monitoring with interactive Slicers,\n"
    "while the Power BI dashboard provides advanced DAX analytics, Star Schema modeling,\n"
    "and 3-month sales forecasting with 95% confidence intervals.",
    16,False,RGBColor(0xBD,0xD7,0xEE),align=PP_ALIGN.CENTER)
deliverables=[("Excel FINAL","Store_Sales_Dashboard_FINAL.xlsx"),
              ("Word Report","BEE_Lab_Report_FINAL.docx"),
              ("Power BI Guide","DAX_Measures.dax + Setup Guide"),
              ("Dataset","retail_sales_final.csv (300 rows)")]
for i,(title,file) in enumerate(deliverables):
    x=0.8+i*3.1
    rect(sl9,x,3.8,2.8,0.6,"2E75B6",title,12,True,WT)
    box(sl9,x,4.45,2.8,0.6,file,10,False,RGBColor(0xBD,0xD7,0xEE),align=PP_ALIGN.CENTER)
box(sl9,0.6,5.5,12,0.5,"Group 4  |  B.Tech First Year  |  BEE Lab  |  April 2026",
    14,False,RGBColor(0xBD,0xD7,0xEE),align=PP_ALIGN.CENTER)
box(sl9,0.6,6.0,12,0.5,"Thank You!",22,True,WT,align=PP_ALIGN.CENTER)

OUT = r"e:\BEE\Section4_Report\BEE_Lab_Presentation_FINAL.pptx"
prs.save(OUT)
print(f"PowerPoint saved -> {OUT}")
print(f"Slides: 9 (Title|Agenda|Problem|Dataset|Excel|PowerBI|DAX|Findings|Conclusion)")
